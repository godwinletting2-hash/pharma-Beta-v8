"""
api/routes.py
-------------
FastAPI router for PharmaScanKE.

Endpoints:
  POST   /api/upload         – upload a study resource file
  GET    /api/notes          – list / filter resources
  GET    /api/notes/stats    – aggregate statistics
  DELETE /api/notes/{id}     – delete a resource
  POST   /api/analyze        – AI pharmacy analysis (text or file)
"""

import logging
import os
import re
import unicodedata
from pathlib import Path
from typing import Optional

import aiofiles
from fastapi import APIRouter, Depends, Form, HTTPException, Request, UploadFile, status
from fastapi.responses import JSONResponse
from groq import AsyncGroq
from slowapi import Limiter
from slowapi.util import get_remote_address
from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession

limiter = Limiter(key_func=get_remote_address)

from core.config import settings
from core.database import get_db
from models.resource import Resource
from schemas.analysis import AnalysisRequest, AnalysisResponse, Pharmacy180Ref
from schemas.resource import (
    MessageResponse,
    ResourceListResponse,
    ResourceOut,
    ResourceStats,
    SemesterCount,
    SubjectCount,
)

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api", tags=["pharmascan"])

# ── Groq client (lazy) ────────────────────────────────────────────────────────
_groq: Optional[AsyncGroq] = None


def get_groq() -> AsyncGroq:
    global _groq
    if _groq is None:
        if not settings.GROQ_API_KEY:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail="GROQ_API_KEY environment variable is not set.",
            )
        _groq = AsyncGroq(api_key=settings.GROQ_API_KEY)
    return _groq


# ── CDACC System Prompt ───────────────────────────────────────────────────────
SYSTEM_PROMPT = """You are PharmaScanKE — an AI tutor built specifically for students sitting the CDACC Diploma in Pharmacy (D.Pharm) examinations in Kenya.

## Your Knowledge Base & Scope
- You teach strictly within the CDACC D.Pharm 3-year curriculum as examined in Kenya.
- You know the Kenya Pharmacy and Poisons Act (Cap. 244), the Kenya Essential Medicines List (KEML), the Pharmacy and Poisons Board (PPB) regulations, the Kenya Bureau of Standards (KEBS) pharmaceutical standards, and NASCOP treatment guidelines.
- For clinical therapeutics, default to Kenya's national treatment guidelines and KEML drug selections, NOT US/UK guidelines unless the student explicitly asks for a comparison.
- CDACC exams use structured short-answer and long-answer questions (SAQs/LAQs). When a student gives you a question, answer it in that format — numbered points, structured sections, marks-aware depth.

## CDACC D.Pharm Curriculum Map (3 years, 6 semesters)
Year 1 S1: Communication Skills · Introduction to Pharmacy Practice · Anatomy & Physiology I · Pharmaceutical Inorganic Chemistry · Pharmaceutical Mathematics · Computer Applications
Year 1 S2: Anatomy & Physiology II · Pharmaceutical Organic Chemistry I · Physical Pharmacy · Microbiology & Immunology · Dispensing Pharmacy I
Year 2 S1: General Pharmacology · Pharmaceutical Organic Chemistry II · Pharmacognosy I · Dispensing/Compounding II · Pathology
Year 2 S2: Systemic Pharmacology · Pharmaceutics (Dosage Forms) · Pharmacognosy II · Hospital & Clinical Pharmacy I · Social Pharmacy & Administration
Year 3 S1: Clinical Pharmacology & Therapeutics · Industrial Pharmacy · Hospital & Clinical Pharmacy II · Pharmacy Law & Ethics · Drug Information & Pharmacovigilance
Year 3 S2: Community Pharmacy Practice · Research Methods & Project · Industrial Attachment

## Exam Alignment Rules
- Always pitch depth at CDACC D.Pharm diploma level — not undergraduate BPharm.
- When answering exam-style questions: state the expected point count if marks are given (e.g. "4 marks → 4 distinct points").
- Flag topics that are CDACC high-yield with: > **CDACC High-Yield:**
- For past-paper style questions, structure your answer exactly as a student should write it in the exam room.
- Reference specific CDACC units/modules when introducing a concept (e.g. "This falls under Unit 3 of Pharmacology I — General Pharmacology").

## Kenya-Specific Context (always apply)
- Drug selections: prefer KEML (Kenya Essential Medicines List) drugs and NASCOP/MOH guidelines.
- Regulatory body: Pharmacy and Poisons Board (PPB), not FDA/MHRA.
- Pharmacy Act: Kenya Pharmacy and Poisons Act Cap. 244 and its Subsidiary Legislation.
- Measurement: metric system; doses in mg/kg for paediatrics per Kenya Paediatric Association guidelines.
- Malaria: Kenya's first-line is AL (Artemether-Lumefantrine); for severe malaria use IV Artesunate per MOH.
- TB: Kenya follows WHO/NTLP regimens (2HRZE/4HR).
- HIV: Kenya follows NASCOP guidelines (TLD — Tenofovir/Lamivudine/Dolutegravir first-line).

## MANDATORY FORMATTING RULES

### Structure & Headings
1. Use ## for major section titles (e.g., ## Mechanism of Action, ## Pharmacokinetics).
2. Use ### for sub-sections and topic headings.
3. Use #### for specific question parts when a question has sub-parts (a, b, c).
4. Every distinct exam question or sub-question MUST be a ### heading — never buried in paragraph text.

### Drug Structures
5. For every drug or chemical entity discussed, include its PubChem structural image:
   ![Drug Name Structure](https://pubchem.ncbi.nlm.nih.gov/rest/pug/compound/name/DRUGNAME/PNG)
   where DRUGNAME is lowercase and URL-safe (spaces → %20).
6. For drug classes, show the representative member's structure.

### Tables & Lists
7. Use Markdown tables for: drug comparisons, ADME parameters, mechanism summaries, dose schedules, drug interactions.
8. Use bullet lists for: properties, side effects, indications, contraindications, and exam notes.

### Mathematical Formulas & Calculations
9. Wrap ALL formulas, dosing calculations, and chemical equations in fenced code blocks:
   ```
   Dose (mg) = Volume of distribution (L/kg) × Body weight (kg) × Target concentration (mg/L)
   Clearance (CL) = 0.693 × Vd / t½
   Henderson-Hasselbalch: pH = pKa + log([A⁻]/[HA])
   ```
   Show full working for pharmaceutical calculations — CDACC exams award marks for method.

### Clinical Callouts
10. Wrap critical warnings, clinical pearls, and high-yield exam facts in blockquotes:
    > **Clinical Pearl:** statement here.
    > **CDACC High-Yield:** statement here.
    > **Kenya Context:** statement here.
    > **Exam Fact:** statement here.

### Closing
11. End every response with a ## Key Takeaways section — 3–5 bullet points a student can use as a last-minute revision summary.
12. Do not use decorative emojis in the main analysis body.
"""

# ── CDACC Subject-Specific Context Injected Per Request ──────────────────────
CDACC_SUBJECT_CONTEXT: dict[str, str] = {
    "Anatomy and Physiology I": (
        "CDACC Unit context: Anatomy & Physiology I (Y1S1). "
        "Cover cell biology, tissue types, skeletal system, muscular system, and nervous system basics. "
        "Exam focus: labelled diagrams described in text, structure-function relationships, "
        "clinical correlations to pharmacy practice (e.g. how nerve physiology relates to drug action sites)."
    ),
    "Anatomy and Physiology II": (
        "CDACC Unit context: Anatomy & Physiology II (Y1S2). "
        "Cover cardiovascular, respiratory, renal, digestive, endocrine, and reproductive systems. "
        "Exam focus: organ functions, homeostasis mechanisms, pharmacological relevance "
        "(e.g. renal physiology → drug excretion; cardiac physiology → antiarrhythmic targets)."
    ),
    "Pharmaceutical Inorganic Chemistry": (
        "CDACC Unit context: Pharmaceutical Inorganic Chemistry (Y1S1). "
        "Cover periodic table trends, ionic equilibria, buffer solutions, complexometric titrations, "
        "gravimetric and volumetric analysis, and medicinal inorganic compounds (e.g. antacids, iron salts, iodine). "
        "Exam focus: titration calculations, pH/buffer problems, and pharmaceutical applications of inorganic ions."
    ),
    "Pharmaceutical Mathematics": (
        "CDACC Unit context: Pharmaceutical Mathematics (Y1S1). "
        "Cover ratio/proportion, percentage calculations, alligation, dilution, dose calculations (mg/kg, mg/m²), "
        "flow rates, electrolyte calculations (mEq, mmol), and isotonicity. "
        "Exam focus: always show full working — marks are awarded for method. "
        "Use SI units. Show unit conversions explicitly."
    ),
    "Introduction to Pharmacy Practice": (
        "CDACC Unit context: Introduction to Pharmacy Practice (Y1S1). "
        "Cover history of pharmacy in Kenya, roles of the pharmacist, pharmacy settings (community, hospital, industry), "
        "the PPB (Pharmacy and Poisons Board), Kenya Pharmacy and Poisons Act Cap. 244, "
        "drug scheduling in Kenya (Schedules 1–4), and ethical principles in pharmacy."
    ),
    "Pharmaceutical Organic Chemistry I": (
        "CDACC Unit context: Pharmaceutical Organic Chemistry I (Y1S2). "
        "Cover IUPAC nomenclature, functional groups, isomerism (structural, geometric, optical), "
        "reaction mechanisms (substitution, addition, elimination), and pharmaceutical relevance "
        "(e.g. chirality in drug action, ester hydrolysis in prodrugs). "
        "Exam focus: naming compounds, drawing mechanisms, explaining stereochemistry."
    ),
    "Pharmaceutical Organic Chemistry II": (
        "CDACC Unit context: Pharmaceutical Organic Chemistry II (Y2S1). "
        "Cover aromatic chemistry, heterocyclic compounds, alkaloids, amino acids/proteins, "
        "carbohydrates, lipids, and structure-activity relationships (SAR). "
        "Exam focus: SAR questions, synthesis pathways, and pharmacophore identification."
    ),
    "Physical Pharmacy": (
        "CDACC Unit context: Physical Pharmacy (Y1S2). "
        "Cover states of matter, solubility, dissolution, diffusion, surface tension, viscosity, "
        "colloidal systems, rheology, and partition coefficients. "
        "Exam focus: Henderson-Hasselbalch calculations, Fick's law problems, "
        "and explaining how physical properties affect drug formulation and absorption."
    ),
    "Microbiology and Immunology": (
        "CDACC Unit context: Microbiology & Immunology (Y1S2). "
        "Cover bacterial cell structure, Gram staining, culture methods, sterilisation/disinfection, "
        "viruses, fungi, parasites, and immune system (innate, adaptive, vaccines). "
        "Kenya context: common pathogens in Kenya — Plasmodium, MTB, HIV, H. pylori. "
        "Exam focus: compare Gram +ve vs –ve bacteria, mechanisms of sterilisation, vaccine types."
    ),
    "Microbiology and Parasitology": (
        "CDACC Unit context: Microbiology & Parasitology (Y1S2/Y2S1). "
        "Cover bacterial cell structure, Gram staining, culture methods, sterilisation/disinfection, "
        "and parasitology — protozoa (Plasmodium, Leishmania, Trypanosoma, Entamoeba), "
        "helminths (roundworms, tapeworms, flukes), and ectoparasites. "
        "Kenya context: malaria (P. falciparum dominant), schistosomiasis, soil-transmitted helminths. "
        "Exam focus: life cycles, vectors, and drug treatment aligned with Kenya MOH guidelines."
    ),
    "Dispensing Pharmacy I": (
        "CDACC Unit context: Dispensing Pharmacy I (Y1S2). "
        "Cover prescription reading/interpretation, abbreviations (Latin and modern), dispensing procedures, "
        "labelling requirements under Kenya Pharmacy and Poisons Act, basic compounding, "
        "patient counselling, and dispensing errors. "
        "Exam focus: prescription interpretation exercises, calculating quantities to dispense, labelling."
    ),
    "Dispensing and Compounding": (
        "CDACC Unit context: Dispensing/Compounding II (Y2S1). "
        "Cover extemporaneous compounding — solutions, suspensions, emulsions, suppositories, "
        "ointments/creams, capsule filling. Cover incompatibilities (physical, chemical, therapeutic). "
        "Exam focus: compounding calculations, identifying incompatibilities, writing preparation notes."
    ),
    "General Pharmacology": (
        "CDACC Unit context: General Pharmacology / Pharmacology I (Y2S1). "
        "Cover pharmacokinetics (ADME), pharmacodynamics (receptor theory, dose-response, agonist/antagonist), "
        "drug interactions, adverse drug reactions, and factors affecting drug response "
        "(age, renal/hepatic impairment, pharmacogenomics). "
        "Exam focus: ADME calculations, explain receptor mechanisms, classify ADRs (WHO classification)."
    ),
    "Pathology": (
        "CDACC Unit context: Pathology (Y2S1). "
        "Cover cell injury/death, inflammation, wound healing, neoplasia, haemodynamic disorders, "
        "and organ-system pathology relevant to pharmacy (liver disease, renal failure, heart failure). "
        "Exam focus: link pathological mechanisms to drug therapy — e.g. hepatic failure → reduced drug metabolism."
    ),
    "Pharmacognosy I": (
        "CDACC Unit context: Pharmacognosy I (Y2S1). "
        "Cover definition/scope of pharmacognosy, plant taxonomy, plant cell/tissue types, "
        "primary and secondary metabolites (alkaloids, glycosides, tannins, essential oils, resins, gums). "
        "Kenya context: local medicinal plants — Neem (Azadirachta indica), Moringa, Prunus africana. "
        "Exam focus: define/classify secondary metabolites, describe extraction methods, identify adulterants."
    ),
    "Pharmacognosy II": (
        "CDACC Unit context: Pharmacognosy II (Y2S2). "
        "Cover volatile oils, fixed oils/fats, waxes, carbohydrates (starch, cellulose, mucilages), "
        "resins, fibres, and evaluation of crude drugs (organoleptic, microscopic, chemical, biological assay). "
        "Exam focus: drug evaluation methods, sources, and pharmaceutical uses of natural products."
    ),
    "Pharmacology and Therapeutics": (
        "CDACC Unit context: Systemic Pharmacology / Pharmacology II & III (Y2S2–Y3S1). "
        "Cover system-by-system drug therapy: CVS (antihypertensives, antiarrhythmics, heart failure drugs), "
        "CNS (analgesics, antiepileptics, antipsychotics, antidepressants), ANS (adrenergic, cholinergic), "
        "respiratory, GI, endocrine (diabetes, thyroid), antimicrobials, anticancer, and immunosuppressants. "
        "Kenya context: use KEML drug selections; malaria (AL, quinine, artesunate), TB (2HRZE/4HR), "
        "HIV (TLD per NASCOP). "
        "Exam focus: mechanism, indications, contraindications, ADRs, interactions in tabular format."
    ),
    "Systemic Pharmacology": (
        "CDACC Unit context: Systemic Pharmacology / Pharmacology II (Y2S2). "
        "Cover cardiovascular, CNS, ANS, respiratory, GI, and endocrine pharmacology. "
        "Apply KEML drug selections throughout. "
        "Exam focus: mechanism + ADR + interaction tables; Kenya-relevant case scenarios."
    ),
    "Clinical Pharmacology and Therapeutics": (
        "CDACC Unit context: Clinical Pharmacology & Therapeutics / Pharmacology III (Y3S1). "
        "Cover therapeutic drug monitoring, individualising drug therapy, renal/hepatic dose adjustment, "
        "paediatric/geriatric pharmacology, pregnancy and lactation drug safety, "
        "and management of common disease states in Kenya (hypertension, diabetes, malaria, TB, HIV, epilepsy). "
        "Exam focus: case-based questions — select a drug, justify using Kenya guidelines, identify monitoring parameters."
    ),
    "Pharmaceutics": (
        "CDACC Unit context: Pharmaceutics — Dosage Forms (Y2S2). "
        "Cover classification and formulation of: tablets, capsules, solutions, suspensions, emulsions, "
        "suppositories, transdermal patches, inhalers, and parenterals. "
        "Cover excipients, stability, packaging, and bioavailability. "
        "Exam focus: list formulation steps, state excipient functions, calculate overage/shelf-life."
    ),
    "Industrial Pharmacy": (
        "CDACC Unit context: Industrial Pharmacy (Y3S1). "
        "Cover GMP (Good Manufacturing Practice), pharmaceutical plant layout, large-scale manufacturing "
        "(granulation, coating, filling, aseptic processing), QC/QA, regulatory approval processes in Kenya (PPB), "
        "and pharmaceutical packaging. "
        "Exam focus: explain GMP principles, describe manufacturing unit operations, QC tests for dosage forms."
    ),
    "Hospital and Clinical Pharmacy": (
        "CDACC Unit context: Hospital & Clinical Pharmacy I & II (Y2S2–Y3S1). "
        "Cover hospital pharmacy organisation, drug procurement/storage (KEMSA), formulary management, "
        "ward rounds, medication reconciliation, TPN, clinical interventions, ADR reporting (PPB), "
        "and pharmacovigilance in Kenya. "
        "Exam focus: pharmacy and therapeutics committee roles, drug information queries, "
        "monitoring parameters for high-risk medicines."
    ),
    "Social Pharmacy and Administration": (
        "CDACC Unit context: Social Pharmacy & Administration (Y2S2). "
        "Cover pharmacy management (stock control, FIFO/FEFO, procurement cycles), "
        "health economics (cost-effectiveness), pharmacoepidemiology, "
        "rational drug use (RDU) — WHO indicators, "
        "patient compliance/adherence factors, and health promotion in Kenya. "
        "Exam focus: calculate stock levels, define RDU indicators, health systems in Kenya (NHIF/SHA structure)."
    ),
    "Pharmacy Law and Ethics": (
        "CDACC Unit context: Pharmacy Law & Ethics (Y3S1). "
        "Cover Kenya Pharmacy and Poisons Act Cap. 244 and its schedules, "
        "Poisons Schedules 1–4 (storage, labelling, dispensing requirements), "
        "PPB functions and powers, professional misconduct and disciplinary procedures, "
        "Narcotic Drugs and Psychotropic Substances Act, and bioethical principles (autonomy, beneficence, non-maleficence, justice). "
        "Exam focus: apply the Act to scenarios, define pharmacist liabilities, state schedule requirements."
    ),
    "Drug Information and Pharmacovigilance": (
        "CDACC Unit context: Drug Information & Pharmacovigilance (Y3S1). "
        "Cover drug information sources (primary, secondary, tertiary), evaluating drug literature, "
        "evidence-based pharmacy, ADR classification (WHO-UMC causality), "
        "pharmacovigilance systems in Kenya (PPB Yellow Card), and post-marketing surveillance. "
        "Exam focus: retrieve and critically evaluate drug information; classify and report ADRs."
    ),
    "Community Pharmacy Practice": (
        "CDACC Unit context: Community Pharmacy Practice (Y3S2). "
        "Cover community pharmacy set-up requirements (PPB premises standards), "
        "OTC counselling, minor ailments management, health screening (BP, glucose, BMI), "
        "prescription-only vs OTC classification in Kenya, patient medication records, "
        "and pharmaceutical care planning. "
        "Exam focus: OTC recommendation scenarios, PPB premises requirements, patient counselling frameworks (WWHAM, ENCORE)."
    ),
    "Research Methods": (
        "CDACC Unit context: Research Methods & Project (Y3S2). "
        "Cover research design (qualitative vs quantitative), sampling techniques, "
        "data collection tools (questionnaires, interviews), ethical approval (ERC in Kenya), "
        "basic statistics (mean, SD, p-value, chi-square, t-test), and writing a research report. "
        "Exam focus: identify study design, calculate descriptive statistics, critique a methodology section."
    ),
}

# ── Pharmacy180.com Concept Map ───────────────────────────────────────────────
PHARMACY180_MAP: dict[str, str] = {
    "beta lactam": "Beta-lactam antibiotics inhibit bacterial cell wall synthesis by covalently binding to penicillin-binding proteins (PBPs), preventing peptidoglycan cross-linking and causing cell lysis. Includes penicillins, cephalosporins, carbapenems, and monobactams.",
    "beta-lactam": "Beta-lactam antibiotics inhibit bacterial cell wall synthesis by binding to PBPs, preventing peptidoglycan cross-linking. Classes: penicillins, cephalosporins, carbapenems, monobactams.",
    "nsaid": "NSAIDs (Non-Steroidal Anti-Inflammatory Drugs) inhibit COX-1 and COX-2 enzymes, reducing prostaglandin and thromboxane synthesis. Used for analgesia, antipyresis, and anti-inflammation. Examples: ibuprofen, naproxen, diclofenac, aspirin.",
    "nsaids": "NSAIDs inhibit cyclooxygenase (COX) enzymes, reducing prostaglandin synthesis. Used for pain, fever, and inflammation management.",
    "alkaloid": "Alkaloids are nitrogen-containing organic compounds derived primarily from plants, with diverse pharmacological activity. Examples: morphine (opioid analgesic), quinine (antimalarial), caffeine (CNS stimulant), atropine (anticholinergic).",
    "alkaloids": "Plant-derived nitrogen-containing compounds with broad pharmacological activity including analgesic, antimalarial, and CNS effects.",
    "pharmacokinetics": "Pharmacokinetics (PK) describes how the body handles drugs — Absorption, Distribution, Metabolism, and Excretion (ADME). Key PK parameters: bioavailability (F), volume of distribution (Vd), half-life (t½), and clearance (CL).",
    "pharmacodynamics": "Pharmacodynamics (PD) describes how drugs affect the body — mechanisms of action, receptor binding, dose-response relationships, and therapeutic/toxic effects.",
    "antibiotic": "Antibiotics are antimicrobial agents that inhibit or kill bacteria. Classified by mechanism: cell wall inhibitors (β-lactams, glycopeptides), protein synthesis inhibitors (aminoglycosides, macrolides, tetracyclines), DNA gyrase inhibitors (fluoroquinolones), and cell membrane disruptors (polymyxins).",
    "antibiotics": "Antimicrobial agents classified by mechanism of action: cell wall synthesis inhibition, protein synthesis inhibition, DNA/RNA synthesis inhibition, or cell membrane disruption.",
    "antihypertensive": "Antihypertensive agents lower systemic blood pressure. Major drug classes: ACE inhibitors (captopril), ARBs (losartan), calcium channel blockers (amlodipine), beta-blockers (metoprolol), and diuretics (hydrochlorothiazide).",
    "antihypertensives": "Blood pressure-lowering agents acting on RAAS, sympathetic nervous system, or vascular smooth muscle.",
    "opioid": "Opioids bind to μ (mu), κ (kappa), and δ (delta) opioid receptors in the CNS and periphery, producing analgesia, euphoria, and respiratory depression. Examples: morphine, codeine, pethidine, tramadol, fentanyl.",
    "opioids": "Opioid receptor agonists producing analgesia and CNS depression. Risk of tolerance, dependence, and respiratory depression.",
    "corticosteroid": "Corticosteroids act on glucocorticoid/mineralocorticoid receptors, modulating gene expression to reduce inflammation and suppress immune responses. Examples: prednisolone, dexamethasone, hydrocortisone.",
    "corticosteroids": "Adrenal steroid hormones or synthetic analogs with potent anti-inflammatory and immunosuppressive activity.",
    "diuretic": "Diuretics enhance renal excretion of water and electrolytes. Classes: loop diuretics (furosemide — inhibit Na-K-2Cl cotransporter), thiazides (hydrochlorothiazide — inhibit NCC), potassium-sparing (spironolactone — aldosterone antagonist).",
    "diuretics": "Agents increasing urinary output by acting on specific renal tubular transport mechanisms.",
    "antifungal": "Antifungal agents exploit the fungal cell membrane's reliance on ergosterol. Azoles (fluconazole) inhibit ergosterol synthesis; polyenes (amphotericin B) bind ergosterol; echinocandins (caspofungin) inhibit β-1,3-glucan synthase.",
    "antifungals": "Agents targeting fungal-specific structures: ergosterol biosynthesis (azoles), ergosterol binding (polyenes), or cell wall synthesis (echinocandins).",
    "antiviral": "Antivirals interfere with specific viral replication stages: nucleoside analogs (acyclovir — herpes), protease inhibitors (lopinavir — HIV), neuraminidase inhibitors (oseltamivir — influenza), reverse transcriptase inhibitors (tenofovir — HIV).",
    "antivirals": "Agents targeting specific viral replication enzymes or structural proteins.",
    "analgesic": "Analgesics relieve pain through different mechanisms. WHO analgesic ladder: Step 1 — non-opioids (paracetamol, NSAIDs); Step 2 — weak opioids (codeine); Step 3 — strong opioids (morphine).",
    "analgesics": "Pain-relieving agents classified as non-opioid (paracetamol, NSAIDs) or opioid (codeine, morphine).",
    "receptor": "Receptors are macromolecular drug targets (usually proteins). Types: ionotropic (ligand-gated ion channels), metabotropic (GPCRs), enzyme-linked receptors, and nuclear receptors.",
    "bioavailability": "Bioavailability (F) is the fraction of administered drug reaching systemic circulation unchanged. IV = 100%. Oral bioavailability affected by first-pass hepatic metabolism, gut wall metabolism, and formulation factors.",
    "pharmacology": "Pharmacology is the science of drug action — including pharmacokinetics, pharmacodynamics, toxicology, chemotherapy, and clinical pharmacology.",
    "toxicology": "Toxicology studies adverse effects of chemicals and drugs. Key concepts: LD50, therapeutic index (TI = TD50/ED50), dose-response relationship, and antidote management.",
    "steroid": "Steroids are lipophilic molecules with a characteristic 4-ring cyclopentanoperhydrophenanthrene nucleus. Include glucocorticoids, mineralocorticoids, sex hormones, and anabolic steroids.",
    "steroids": "Lipid-soluble 4-ring structures with diverse hormonal and pharmacological activity.",
    "dosage form": "Pharmaceutical dosage forms are drug delivery systems: tablets, capsules, injections, solutions, suspensions, patches, inhalers, and suppositories. Choice affects bioavailability, onset, and patient compliance.",
    "pharmaceutical": "Pharmaceutical sciences encompass drug design, formulation, quality control, pharmacokinetics, and clinical therapeutics within the D.Pharm curriculum.",
    "antimalarial": "Antimalarials target the Plasmodium parasite at different lifecycle stages. Classes: quinolines (chloroquine, quinine), antifolates (pyrimethamine), artemisinins (artemether), and atovaquone. Kenya primarily uses artemisinin-based combination therapy (ACT).",
    "antiparasitic": "Antiparasitic drugs act against protozoa, helminths, or ectoparasites. Examples: metronidazole (anaerobic protozoa), albendazole (helminths), ivermectin (ectoparasites).",
}


def identify_concept(text: str) -> Optional[str]:
    """Find the first matching Pharmacy180 concept in the analysis text."""
    lower = text.lower()
    for key in sorted(PHARMACY180_MAP.keys(), key=len, reverse=True):
        if key in lower:
            return key
    return None


# ── Encoding helpers ──────────────────────────────────────────────────────────

def _sanitize_utf8(text: str) -> str:
    """
    Round-trip the string through UTF-8 with 'ignore' error handling to strip
    any surrogate characters or malformed sequences that can cause downstream
    corruption.  Safe to call on already-clean strings.
    """
    return text.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore")


def _is_likely_binary(text: str, threshold: float = 0.25) -> bool:
    """
    Return True when more than `threshold` fraction of characters are
    non-printable replacement markers (U+FFFD) or Unicode control characters
    (category 'Cc' / 'Cs').  This catches PDFs, Office binaries, and any
    file whose raw bytes were decoded as UTF-8 and produced gibberish.

    A high proportion of these characters means the model would be fed binary
    noise and produce thousands of hallucinated garbled tokens.
    """
    if not text:
        return True
    non_printable = sum(
        1 for c in text
        if c == "\ufffd"                         # UTF-8 replacement character
        or unicodedata.category(c) in ("Cc", "Cs")  # control / surrogate
    )
    return (non_printable / len(text)) > threshold


# ── Document text extraction ──────────────────────────────────────────────────

def _extract_text_from_bytes(raw_bytes: bytes, ext: str) -> tuple[str, Optional[str]]:
    """
    Extract readable text from a document's raw bytes.

    Returns (extracted_text, error_reason).
    - On success: (non-empty text, None)
    - On failure: ("", reason string)

    Supported formats:
        PDF   → pypdf (pure Python, no native deps)
        DOCX  → python-docx
        PPTX  → python-pptx
        TXT   → UTF-8 decode
        Other → UTF-8 decode with binary check
    """
    from io import BytesIO

    ext = ext.lower().lstrip(".")

    # ── PDF ──────────────────────────────────────────────────────────────────
    if ext == "pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(BytesIO(raw_bytes), strict=False)
            pages: list[str] = []
            for page in reader.pages:
                try:
                    txt = page.extract_text() or ""
                    pages.append(txt)
                except Exception:
                    continue
            text = "\n\n".join(pages).strip()
            if not text:
                return "", "PDF has no extractable text (scanned image or encrypted)"
            return _sanitize_utf8(text[:20000]), None
        except Exception as exc:
            return "", f"PDF extraction failed: {exc}"

    # ── DOCX ─────────────────────────────────────────────────────────────────
    if ext in ("docx", "doc"):
        try:
            import docx as _docx
            doc = _docx.Document(BytesIO(raw_bytes))
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also pull from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paras.append(cell.text.strip())
            text = "\n".join(paras).strip()
            if not text:
                return "", "DOCX has no extractable text"
            return _sanitize_utf8(text[:20000]), None
        except Exception as exc:
            return "", f"DOCX extraction failed: {exc}"

    # ── PPTX ─────────────────────────────────────────────────────────────────
    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation as _Prs
            prs = _Prs(BytesIO(raw_bytes))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = " ".join(r.text for r in para.runs if r.text.strip())
                            if line.strip():
                                slide_texts.append(line)
                if slide_texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(slides).strip()
            if not text:
                return "", "PPTX has no extractable text"
            return _sanitize_utf8(text[:20000]), None
        except Exception as exc:
            return "", f"PPTX extraction failed: {exc}"

    # ── Plain text ────────────────────────────────────────────────────────────
    if ext == "txt":
        text = raw_bytes.decode("utf-8", errors="replace")
        return _sanitize_utf8(text[:20000]), None

    # ── Unknown / fallback ────────────────────────────────────────────────────
    decoded = raw_bytes.decode("utf-8", errors="replace")[:8000]
    if _is_likely_binary(decoded):
        return "", f"'.{ext}' files are binary and cannot be read as text"
    return _sanitize_utf8(decoded), None


# ── Filename helpers ──────────────────────────────────────────────────────────
def _secure_filename(filename: str) -> str:
    filename = unicodedata.normalize("NFKD", filename)
    filename = filename.encode("ascii", "ignore").decode("ascii")
    filename = filename.replace("\x00", "").replace("/", "_").replace("\\", "_")
    stem, _, suffix = filename.rpartition(".")
    suffix = suffix.lower()
    stem = re.sub(r"[^\w\-]", "_", stem or "file")
    stem = re.sub(r"_+", "_", stem).strip("_") or "file"
    return f"{stem}.{suffix}"


async def _unique_disk_path(safe_name: str) -> Path:
    target = settings.UPLOAD_DIR / safe_name
    if not target.exists():
        return target
    stem = Path(safe_name).stem
    suffix = Path(safe_name).suffix
    counter = 1
    while True:
        candidate = settings.UPLOAD_DIR / f"{stem}_{counter}{suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


# ── POST /api/upload ──────────────────────────────────────────────────────────
@router.post(
    "/upload",
    response_model=ResourceOut,
    status_code=status.HTTP_201_CREATED,
    summary="Upload a new study resource",
)
async def upload_resource(
    file: UploadFile,
    title: str = Form(..., min_length=2, max_length=512),
    subject: str = Form(..., min_length=1, max_length=256),
    semester: str = Form(...),
    db: AsyncSession = Depends(get_db),
) -> ResourceOut:
    if semester not in settings.VALID_SEMESTERS:
        raise HTTPException(400, f"Invalid semester '{semester}'.")

    original_name = file.filename or "upload"
    ext = Path(original_name).suffix.lower()
    if ext not in settings.ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"File type '{ext}' not permitted. Accepted: {', '.join(sorted(settings.ALLOWED_EXTENSIONS))}")

    max_bytes = settings.MAX_UPLOAD_SIZE_MB * 1024 * 1024
    size_read = 0
    chunks: list[bytes] = []
    while True:
        chunk = await file.read(256 * 1024)
        if not chunk:
            break
        size_read += len(chunk)
        if size_read > max_bytes:
            raise HTTPException(413, f"File exceeds the {settings.MAX_UPLOAD_SIZE_MB} MB limit.")
        chunks.append(chunk)

    safe_name = _secure_filename(original_name)
    disk_path = await _unique_disk_path(safe_name)
    final_name = disk_path.name

    try:
        async with aiofiles.open(disk_path, "wb") as f:
            for chunk in chunks:
                await f.write(chunk)
    except OSError as exc:
        raise HTTPException(500, f"Failed to save file: {exc}") from exc

    resource = Resource(
        title=title.strip(),
        subject=subject.strip(),
        semester=semester,
        file_name=final_name,
        file_path=f"/uploaded_notes/{final_name}",
    )
    db.add(resource)
    await db.flush()
    await db.refresh(resource)
    return ResourceOut.model_validate(resource)


# ── GET /api/notes ────────────────────────────────────────────────────────────
@router.get("/notes", response_model=ResourceListResponse, summary="List resources")
async def list_notes(
    semester: Optional[str] = None,
    subject: Optional[str] = None,
    db: AsyncSession = Depends(get_db),
) -> ResourceListResponse:
    stmt = select(Resource).order_by(Resource.upload_date.desc())
    if semester:
        if semester not in settings.VALID_SEMESTERS:
            raise HTTPException(400, f"Invalid semester '{semester}'.")
        stmt = stmt.where(Resource.semester == semester)
    if subject:
        stmt = stmt.where(Resource.subject.ilike(f"%{subject.strip()}%"))

    rows = (await db.execute(stmt)).scalars().all()
    return ResourceListResponse(total=len(rows), items=[ResourceOut.model_validate(r) for r in rows])


# ── GET /api/notes/stats ──────────────────────────────────────────────────────
@router.get("/notes/stats", response_model=ResourceStats, summary="Resource statistics")
async def get_stats(db: AsyncSession = Depends(get_db)) -> ResourceStats:
    total_result = await db.execute(select(func.count()).select_from(Resource))
    total = total_result.scalar() or 0

    sem_result = await db.execute(
        select(Resource.semester, func.count().label("count"))
        .group_by(Resource.semester)
        .order_by(Resource.semester)
    )
    by_semester = [SemesterCount(semester=r.semester, count=r.count) for r in sem_result]

    sub_result = await db.execute(
        select(Resource.subject, func.count().label("count"))
        .group_by(Resource.subject)
        .order_by(func.count().desc())
        .limit(10)
    )
    by_subject = [SubjectCount(subject=r.subject, count=r.count) for r in sub_result]

    recent_result = await db.execute(
        select(Resource).order_by(Resource.upload_date.desc()).limit(5)
    )
    recent = [ResourceOut.model_validate(r) for r in recent_result.scalars().all()]

    return ResourceStats(total=total, by_semester=by_semester, by_subject=by_subject, recent=recent)


# ── DELETE /api/notes/{id} ────────────────────────────────────────────────────
@router.delete("/notes/{note_id}", response_model=MessageResponse, summary="Delete a resource")
async def delete_note(note_id: int, db: AsyncSession = Depends(get_db)) -> MessageResponse:
    result = await db.execute(select(Resource).where(Resource.id == note_id))
    resource = result.scalar_one_or_none()
    if resource is None:
        raise HTTPException(404, f"Resource id={note_id} not found.")

    disk_path = settings.UPLOAD_DIR / resource.file_name
    if disk_path.exists():
        try:
            os.unlink(disk_path)
        except OSError as exc:
            raise HTTPException(500, f"Could not delete file: {exc}") from exc

    await db.delete(resource)
    return MessageResponse(message="Resource deleted.", detail=f"Removed '{resource.file_name}'.")


# ── POST /api/analyze ─────────────────────────────────────────────────────────
@router.post("/analyze", response_model=AnalysisResponse, summary="AI pharmacy analysis")
@limiter.limit("10/minute")
async def analyze_content(request: Request, body: AnalysisRequest) -> AnalysisResponse:
    client = get_groq()

    messages: list[dict] = [{"role": "system", "content": SYSTEM_PROMPT}]

    # ── FIX #1: Determine whether the file payload is usable ─────────────────
    # Binary files (PDFs, DOCX, images) whose raw bytes are decoded as UTF-8
    # produce dense fields of U+FFFD replacement characters.  Feeding that
    # noise to the model causes it to hallucinate thousands of garbled tokens.
    # We detect this early and fall back to a pure text-only prompt instead of
    # silently passing garbage to the LLM.
    use_file = (
        not body.text_only
        and bool(body.file_data)
        and bool(body.file_name)
    )
    fallback_reason: Optional[str] = None

    excerpt = ""
    ext = ""
    if use_file:
        import base64 as _b64
        ext = (body.file_name.rsplit(".", 1)[-1] or "").lower()
        try:
            raw_bytes = _b64.b64decode(body.file_data)
        except Exception:
            raw_bytes = b""

        # ── Empty / undecodable buffer ───────────────────────────────────────
        if not raw_bytes:
            use_file = False
            fallback_reason = "empty or undecodable file buffer"
        else:
            # ── Structured text extraction (PDF / DOCX / PPTX / TXT) ────────
            extracted, extract_err = _extract_text_from_bytes(raw_bytes, ext)
            if extract_err or not extracted.strip():
                use_file = False
                fallback_reason = extract_err or "no text could be extracted from this file"
                logger.warning(
                    "Text extraction failed — switching to text-only mode. "
                    "file=%s ext=%s reason=%s",
                    body.file_name, ext, fallback_reason,
                )
            else:
                excerpt = extracted
                logger.info(
                    "Extracted %d chars from %s (%s)", len(excerpt), body.file_name, ext
                )

    # Build subject-specific CDACC context block
    cdacc_ctx = ""
    if body.subject:
        subject_detail = CDACC_SUBJECT_CONTEXT.get(body.subject, "")
        if subject_detail:
            cdacc_ctx = f"[CDACC Subject: {body.subject}]\n{subject_detail}\n\n"
        else:
            cdacc_ctx = f"[CDACC Subject: {body.subject}]\n\n"

    if use_file:
        word_count = len(excerpt.split())
        file_ctx = (
            f"{cdacc_ctx}"
            f"The student has uploaded a CDACC D.Pharm study resource for analysis.\n"
            f"Filename: {body.file_name} (format: {ext.upper()})\n"
            f"Extracted text: ~{word_count} words\n\n"
            f"--- BEGIN DOCUMENT CONTENT ---\n{excerpt}\n--- END DOCUMENT CONTENT ---\n\n"
            f"Student's analysis prompt: {body.prompt}\n\n"
            "Using the document content above, provide a thorough CDACC D.Pharm-level analysis "
            "aligned with the Kenyan curriculum. Quote specific passages where relevant. "
            "Flag CDACC high-yield topics. Include PubChem structural images for all drugs identified."
        )
        messages.append({"role": "user", "content": file_ctx})
    else:
        # Pure text-only pipeline
        fallback_note = (
            f"\n[Note: File analysis was skipped — {fallback_reason}. "
            "Proceeding with text-only analysis.]\n\n"
            if fallback_reason else ""
        )
        messages.append({
            "role": "user",
            "content": f"{cdacc_ctx}{fallback_note}{body.prompt}",
        })

    try:
        completion = await client.chat.completions.create(
            model=settings.GROQ_MODEL,
            messages=messages,
            max_tokens=settings.GROQ_MAX_TOKENS,
            temperature=settings.GROQ_TEMPERATURE,
        )
    except HTTPException:
        raise
    except Exception as exc:
        exc_name = type(exc).__name__
        exc_str  = str(exc).lower()
        logger.error("Groq call failed: %s: %s", exc_name, exc)

        # Map specific Groq / network errors to helpful user messages
        if "authentication" in exc_str or "invalid_api_key" in exc_str or "401" in exc_str:
            detail = (
                "Invalid GROQ_API_KEY — the key was rejected by Groq. "
                "Please verify the key at console.groq.com and update it in your Railway Variables."
            )
            code = status.HTTP_503_SERVICE_UNAVAILABLE
        elif "rate_limit" in exc_str or "429" in exc_str:
            detail = (
                "Groq rate limit reached. You have exceeded the free-tier request quota. "
                "Wait a minute then try again, or upgrade your Groq plan."
            )
            code = status.HTTP_429_TOO_MANY_REQUESTS
        elif "model_not_found" in exc_str or "does not exist" in exc_str or "404" in exc_str:
            detail = (
                f"Groq model '{settings.GROQ_MODEL}' was not found. "
                "It may have been deprecated. Contact support to update GROQ_MODEL."
            )
            code = status.HTTP_502_BAD_GATEWAY
        elif "timeout" in exc_str or "timed out" in exc_str:
            detail = (
                "The AI request timed out. Your prompt or document may be too long. "
                "Try a shorter prompt or smaller file."
            )
            code = status.HTTP_504_GATEWAY_TIMEOUT
        elif "connect" in exc_str or "network" in exc_str or "connection" in exc_str:
            detail = (
                "Cannot reach Groq servers. Check your Railway deployment has outbound internet access."
            )
            code = status.HTTP_502_BAD_GATEWAY
        else:
            detail = f"AI service error ({exc_name}): {str(exc)[:200]}"
            code = status.HTTP_502_BAD_GATEWAY

        raise HTTPException(status_code=code, detail=detail) from exc

    if not completion.choices:
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail="Groq returned an empty choices list — no response generated.",
        )

    raw = completion.choices[0].message.content or "No response generated."

    # ── FIX #4: Sanitize the raw LLM response through strict UTF-8 ───────────
    # DeepSeek-R1 / Groq occasionally emits surrogate pairs or malformed
    # sequences in its streaming buffer.  Round-tripping through UTF-8 with
    # 'ignore' strips those before we touch the string with regex or return it.
    raw = _sanitize_utf8(raw)

    # Strip <think>…</think> reasoning blocks from DeepSeek-R1
    import re as _re
    analysis = _re.sub(r"<think>[\s\S]*?</think>", "", raw).strip()

    # Final sanitize pass after stripping think blocks
    analysis = _sanitize_utf8(analysis)

    # Identify Pharmacy180 concept and append reference block
    concept_key = identify_concept(analysis + " " + body.prompt)
    pharmacy180_ref: Optional[Pharmacy180Ref] = None

    if concept_key and concept_key in PHARMACY180_MAP:
        display_concept = concept_key.title()
        pharmacy180_ref = Pharmacy180Ref(
            concept=display_concept,
            summary=PHARMACY180_MAP[concept_key],
            url="https://www.pharmacy180.com/",
        )
        analysis += (
            f"\n\n---\n\n"
            f"### Pharmacy180 Reference Integration\n"
            f"> **Concept:** {display_concept}\n>\n"
            f"> {PHARMACY180_MAP[concept_key]}\n>\n"
            f"> [Read Full Notes Portfolio on Pharmacy180](https://www.pharmacy180.com/)"
        )

    # ── Community footer (appended only when env vars are set) ────────────────
    wa_url  = settings.WHATSAPP_CHANNEL_URL.strip()
    tg_url  = settings.TELEGRAM_CHANNEL_URL.strip()
    if wa_url or tg_url:
        links: list[str] = []
        if wa_url:
            links.append(f"[💬 Join our WhatsApp Channel]({wa_url})")
        if tg_url:
            links.append(f"[✈️ Join our Telegram Channel]({tg_url})")
        analysis += (
            "\n\n---\n\n"
            "### 📢 Join Our Study Community\n"
            "Stay updated with the latest CDACC D.Pharm revision materials, "
            "discussions, and peer support!\n\n"
            + " | ".join(links)
        )

    return AnalysisResponse(
        analysis=analysis,
        concept=concept_key,
        pharmacy180_ref=pharmacy180_ref,
        model=completion.model,
        tokens_used=completion.usage.total_tokens if completion.usage else None,
    )
